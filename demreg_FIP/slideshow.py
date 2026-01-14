import os
import shutil
from pptx import Presentation
from pptx.util import Inches
import glob
import re
from collections import defaultdict

base_dir = os.getcwd()  # Expected: C:\Users\domor\gazelle\demreg_FIP
results_dir = os.path.join(base_dir, "Results")
intensity_folder = os.path.join(base_dir, "intensity_maps")
slides_folder = os.path.join(base_dir, "slides")
overlayed_folder = r"C:\Users\domor\gazelle\demreg_FIP\overlayed_images"
output_pptx = os.path.join(slides_folder, "presentation.pptx")

line_databases = {
    "sis": ['si_10_258.37', 's_10_264.23', 'SiX_SX'],
    "sar": ['s_11_188.68', 'ar_11_188.81', 'SXI_ArXI'],
    "CaAr": ['ca_14_193.87', 'ar_14_194.40', 'CaXIV_ArXIV'],
    "FeS": ['fe_16_262.98', 's_13_256.69', 'FeXVI_SXIII'],
}

def transform_code(code):
    """Transform a code from, e.g., "ca_14_193.87" into "ca14193_87"."""
    return code.replace("_", "").replace(".", "_")

# --- Copy Files from Results/Intensity to Slides Folder
error_log = []

# Loop through all folders in Results that start with "eis_"
for folder in os.listdir(results_dir):
    folder_path = os.path.join(results_dir, folder)
    if os.path.isdir(folder_path) and folder.startswith("eis_"):
        # Extract timestamp from folder name, e.g. "eis_20140202_141952" -> "20140202_141952"
        timestamp = folder[len("eis_"):]
        # Convert timestamp for intensity maps: "20140202_141952" -> "2014_02_02__14_19_52"
        date_part = timestamp[:8]
        time_part = timestamp[9:]  # skip underscore at index 8
        converted_timestamp = f"{date_part[:4]}_{date_part[4:6]}_{date_part[6:]}__{time_part[:2]}_{time_part[2:4]}_{time_part[4:]}"
        # Copy a single overlayed image for this timestamp
        overlayed_pattern = os.path.join(overlayed_folder, f"overlayed_{converted_timestamp}_*.png")
        matching_overlayed_files = glob.glob(overlayed_pattern)
        if matching_overlayed_files:
            # Select the first matching overlayed image
            selected_overlayed_image = matching_overlayed_files[0]
            dst_overlayed = os.path.join(slides_folder, os.path.basename(selected_overlayed_image))
            # Copy the overlayed image if it exists
            if os.path.exists(selected_overlayed_image):
                try:
                    shutil.copy(selected_overlayed_image, dst_overlayed)
                    print(f"Copied overlayed image: {selected_overlayed_image}")
                except Exception as e:
                    error_log.append(f"Error copying overlayed image '{selected_overlayed_image}': {e}")
            else:
                error_log.append(f"Overlayed image not found: {selected_overlayed_image}")
        else:
            error_log.append(f"No overlayed image found for timestamp {converted_timestamp}")

        # Loop over each channel in the database
        for channel in line_databases:
            db_entry = line_databases[channel]
            code1 = db_entry[0]  # first intensity code
            code2 = db_entry[1]  # second intensity code
            expected_code1 = transform_code(code1)
            expected_code2 = transform_code(code2)
            # Build the EIS filename (e.g. eis_20140202_141952_CaAr_map.png)
            eis_filename = f"eis_{timestamp}_{channel}_map.png"
            src_eis = os.path.join(folder_path, eis_filename)
            dst_eis = os.path.join(slides_folder, eis_filename)
            
            # Build the intensity map filenames (e.g. 2014_02_02__14_19_52_ca14193_87.png)
            intensity_filename1 = f"{converted_timestamp}_{expected_code1}.png"
            intensity_filename2 = f"{converted_timestamp}_{expected_code2}.png"
            src_int1 = os.path.join(intensity_folder, intensity_filename1)
            src_int2 = os.path.join(intensity_folder, intensity_filename2)
            dst_int1 = os.path.join(slides_folder, intensity_filename1)
            dst_int2 = os.path.join(slides_folder, intensity_filename2)
            
            # Copy the EIS file if it exists
            if os.path.exists(src_eis):
                try:
                    shutil.copy(src_eis, dst_eis)
                except Exception as e:
                    error_log.append(f"Error copying EIS file '{src_eis}' for channel '{channel}': {e}")
            else:
                error_log.append(f"EIS file not found: {src_eis}")
            
            # Copy the first intensity map if it exists
            if os.path.exists(src_int1):
                try:
                    shutil.copy(src_int1, dst_int1)
                except Exception as e:
                    error_log.append(f"Error copying intensity file '{src_int1}' for channel '{channel}': {e}")
            else:
                error_log.append(f"Intensity file not found: {src_int1}")
            
            # Copy the second intensity map if it exists
            if os.path.exists(src_int2):
                try:
                    shutil.copy(src_int2, dst_int2)
                except Exception as e:
                    error_log.append(f"Error copying intensity file '{src_int2}' for channel '{channel}': {e}")
            else:
                error_log.append(f"Intensity file not found: {src_int2}")
            
            # Print summary for this channel
            if os.path.exists(src_eis) and os.path.exists(src_int1) and os.path.exists(src_int2):
                print(f"Copied all files for folder '{folder}' and channel '{channel}'.")
            else:
                print(f"Some files missing for folder '{folder}' and channel '{channel}'")


# Process files in the slides folder
for filename in os.listdir(slides_folder):
    if filename.endswith(".png"):
        old_path = os.path.join(slides_folder, filename)  # Full path of old file
        new_name = None  # Initialize new name

        # Check if it's a composition map
        match_comp = re.match(r"eis_(\d{8})_(\d{6})_([A-Za-z]+)_map\.png", filename)
        if match_comp:
            date, time, channel = match_comp.groups()
            new_name = f"composition_map_{date[:4]}_{date[4:6]}_{date[6:]}__{time[:2]}_{time[2:4]}_{time[4:]}_{channel}.png"

        # Check if it's an intensity map
        match_int = re.match(r"(\d{4}_\d{2}_\d{2}__\d{2}_\d{2}_\d{2})_(.+)\.png", filename)
        if match_int:
            timestamp, element = match_int.groups()
            new_name = f"intensity_map_{timestamp}_{element}.png"

        # If a new name was generated, rename the file
        if new_name:
            new_path = os.path.join(slides_folder, new_name)
            
            # Check if the file already exists
            if os.path.exists(new_path):
                print(f"Skipping rename: {filename} -> {new_name} (file already exists)")
            else:
                os.rename(old_path, new_path)  # Rename only if no conflict
                print(f"Renamed: {filename} -> {new_name}")


error_log_path = os.path.join(slides_folder, "error_log.txt")
with open(error_log_path, "w", encoding="utf-8") as log_file:
    for err in error_log:
        log_file.write(err + "\n")
        
print("File processing complete.")

# --- Create a PowerPoint Presentation from the Slides Folder

channel_order = ["sis", "sar", "CaAr", "FeS"]
# Expected intensity mappings per composition channel
line_databases = {
    "sis": ['si10258_37', 's10264_23'],
    "sar": ['s11188_68', 'ar11188_81'],
    "CaAr": ['ca14193_87', 'ar14194_40'],
    "FeS": ['fe16262_98', 's13256_69'],
}

# Dictionary to store grouped files
grouped_files = defaultdict(list)

# Scan through files in slides_folder
for filename in os.listdir(slides_folder):
    if filename.endswith(".png"):
        # Extract timestamp (YYYY_MM_DD__HH_MM_SS) from overlayed, composition, or intensity files
        match = re.match(r"(?:overlayed_|composition_map_|intensity_map_)?(\d{4}_\d{2}_\d{2}__\d{2}_\d{2}_\d{2})", filename)
        if match:
            timestamp = match.group(1)
            grouped_files[timestamp].append(filename)


# Function to sort files within each group
def sort_group(files):
    overlayed = []
    composition_maps = []
    intensity_maps = defaultdict(list)
    # First, categorize files
    for file in files:
        if file.startswith("overlayed_"):
            overlayed.append(file)
        elif file.startswith("composition_map_"):
            composition_maps.append(file)
        elif file.startswith("intensity_map_"):
            match = re.match(r"intensity_map_(\d{4}_\d{2}_\d{2}__\d{2}_\d{2}_\d{2})_(.+)\.png", file)
            if match:
                timestamp, element = match.groups()
                intensity_maps[(timestamp, element)].append(file)  # Store by timestamp + element
    # Sort composition maps by predefined channel order
    composition_maps.sort(
        key=lambda f: channel_order.index(f.split("_")[-1].split(".")[0]) if f.split("_")[-1].split(".")[0] in channel_order else len(channel_order)
    )
    # Build final sorted list
    sorted_files = overlayed.copy()
    for comp in composition_maps:
        sorted_files.append(comp)
        # Extract timestamp and channel from the composition map
        match = re.match(r"composition_map_(\d{4}_\d{2}_\d{2}__\d{2}_\d{2}_\d{2})_(.+)\.png", comp)
        if match:
            timestamp, channel = match.groups()
            expected_intensity_codes = line_databases.get(channel, [])
            # Append intensity maps that match both timestamp and expected intensity code
            for intensity_code in expected_intensity_codes:
                matched_intensity = intensity_maps.get((timestamp, intensity_code), [])
                matched_intensity.sort()
                sorted_files.extend(matched_intensity)
    return sorted_files


# Print results with sorted order
for timestamp, files in grouped_files.items():
    sorted_files = sort_group(files)
    print(f"Group: {timestamp} ({len(sorted_files)} files)")
    for f in sorted_files:
        print("  ", f)
    print()  # Newline for better readability


# Sort timestamps from oldest to newest
sorted_timestamps = sorted(grouped_files.keys())

# Create a new PowerPoint presentation
prs = Presentation()
slide_layout = prs.slide_layouts[6]  # Blank slide layout

# Add images to PowerPoint in correct order
for timestamp in sorted_timestamps:
    sorted_files = sort_group(grouped_files[timestamp])
    for filename in sorted_files:
        img_path = os.path.join(slides_folder, filename)
        slide = prs.slides.add_slide(slide_layout)
        try:
            slide.shapes.add_picture(img_path, left=0, top=0, width=prs.slide_width, height=prs.slide_height)
        except Exception as e: #if deleted or moved file
            print(f"Error adding image {filename} to PowerPoint: {e}")

# Save PowerPoint file
prs.save(output_pptx)
print(f"PowerPoint saved as: {output_pptx}")

